import os
from flask import Flask, request, render_template, send_file, send_from_directory
from werkzeug.utils import secure_filename
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from PIL import Image

app = Flask(__name__)
app.config['UPLOAD_FOLDER'] = 'static/uploads'
app.config['OUTPUT_FOLDER'] = 'output'

os.makedirs(app.config['UPLOAD_FOLDER'], exist_ok=True)
os.makedirs(app.config['OUTPUT_FOLDER'], exist_ok=True)

def create_ppt(image_path, num_slides, slide_texts=None, text_color="#FFFFFF"):
    prs = Presentation()
    blank_layout = prs.slide_layouts[6]
    rgb = RGBColor.from_string(text_color.strip("#"))

    for i in range(num_slides):
        slide = prs.slides.add_slide(blank_layout)
        slide.shapes.add_picture(image_path, 0, 0, width=prs.slide_width, height=prs.slide_height)
        if slide_texts and i < len(slide_texts):
            textbox = slide.shapes.add_textbox(Inches(1), Inches(6), Inches(8), Inches(1))
            tf = textbox.text_frame
            p = tf.add_paragraph()
            p.text = slide_texts[i]
            p.font.size = Pt(28)
            p.font.color.rgb = rgb

    output_path = os.path.join(app.config['OUTPUT_FOLDER'], 'generated_presentation.pptx')
    prs.save(output_path)
    return output_path

@app.route("/", methods=["GET"])
def index():
    return render_template("index.html")


@app.route("/preview", methods=["POST"])
def preview():
    image = request.files["image"]
    num_slides = int(request.form["num_slides"])
    slide_texts = request.form.get("slide_texts", "").split("\n")
    text_color = request.form.get("text_color", "#FFFFFF")

    filename = secure_filename(image.filename)
    image_path = os.path.join(app.config['UPLOAD_FOLDER'], filename)
    image.save(image_path)

    return render_template("preview.html", image=filename, num_slides=num_slides, slide_texts=slide_texts, text_color=text_color)


@app.route("/download", methods=["POST"])
def download():
    image = request.form["image"]
    num_slides = int(request.form["num_slides"])
    slide_texts = request.form.get("slide_texts", "").split("\n")
    text_color = request.form.get("text_color", "#FFFFFF")
    image_path = os.path.join(app.config['UPLOAD_FOLDER'], image)
    ppt_path = create_ppt(image_path, num_slides, slide_texts, text_color)
    return send_file(ppt_path, as_attachment=True)

# Serve manifest and service worker
@app.route('/manifest.json')
def manifest():
    return send_from_directory('.', 'manifest.json', mimetype='application/manifest+json')

@app.route('/sw.js')
def sw():
    return send_from_directory('.', 'sw.js', mimetype='application/javascript')

if __name__ == "__main__":
    app.run(debug=True)
